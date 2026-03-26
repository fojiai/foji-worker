from pptx import Presentation
from pptx.util import Pt


def extract(file_path: str) -> str:
    prs = Presentation(file_path)
    slide_texts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        texts.append(text)
        if texts:
            slide_texts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slide_texts)
